# Standard library imports
import os
import re
import xml.etree.ElementTree as ET
import zipfile
from io import BytesIO
# Third-party imports
import openai
import PyPDF2
import textract
import tiktoken
from bs4 import BeautifulSoup
from docx import Document
from lxml import etree
from PIL import Image
from pptx import Presentation
# Local imports
from src import config

openai.api_key = os.getenv('openai_api_key')

def parse_file(attachment, file_content):
  file_text = ""

  # Process different file types
  if attachment.content_type == "text/plain":
      file_text = file_content.decode("utf-8")
  elif attachment.content_type == "application/pdf":
      with BytesIO(file_content) as pdf_file:
          pdf_reader = PyPDF2.PdfReader(pdf_file)
          for page_num in range(len(pdf_reader.pages)):
              file_text += pdf_reader.pages[page_num].extract_text()
  elif attachment.content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
      with BytesIO(file_content) as docx_file:
          doc = Document(docx_file)
          file_text = "\n".join(para.text for para in doc.paragraphs)
  elif attachment.content_type == "application/vnd.ms-powerpoint":
      with BytesIO(file_content) as ppt_file:
          ppt = Presentation(ppt_file)
          for slide in ppt.slides:
              for shape in slide.shapes:
                  if hasattr(shape, "text"):
                      file_text += shape.text + "\n"
  elif attachment.content_type.startswith("image/"):
      with BytesIO(file_content) as img_file:
          img = Image.open(img_file)
          file_text = f"\n\nImage file: {attachment.filename}\n"
  # Add support for TeX files
  elif attachment.filename.endswith((".tex", ".ltx")):
      file_text = file_content.decode("utf-8")
  # Add support for XML files
  elif attachment.filename.endswith(".xml"):
      tree = ET.ElementTree(ET.fromstring(file_content))
      file_text = etree.tostring(tree, encoding='utf-8', method='xml', pretty_print=True).decode("utf-8")
  # Add support for HTML files
  elif attachment.filename.endswith(".html") or attachment.content_type.startswith("text/html"):
      soup = BeautifulSoup(file_content, "html.parser")
      file_text = soup.get_text()
  # Add support for Python files
  elif attachment.filename.endswith(".py"):
      file_text = file_content.decode("utf-8")
  # Support for other plain text file types
  elif attachment.content_type.startswith("text/"):
      file_text = file_content.decode("utf-8")
  # Support for other file types using textract
  else:
      try:
          file_text = textract.process(BytesIO(file_content)).decode("utf-8")
      except:
          file_text = ""

  return file_text

def num_tokens_from_messages(messages, model="gpt-3.5-turbo-0301"):
    """Returns the number of tokens used by a list of messages."""
    try:
        encoding = tiktoken.encoding_for_model(model)
    except KeyError:
        print("Warning: model not found. Using cl100k_base encoding.")
        encoding = tiktoken.get_encoding("cl100k_base")
    if model == "gpt-3.5-turbo":
        print("Warning: gpt-3.5-turbo may change over time. Returning num tokens assuming gpt-3.5-turbo-0301.")
        return num_tokens_from_messages(messages, model="gpt-3.5-turbo-0301")
    elif model == "gpt-4":
        print("Warning: gpt-4 may change over time. Returning num tokens assuming gpt-4-0314.")
        return num_tokens_from_messages(messages, model="gpt-4-0314")
    elif model == "gpt-3.5-turbo-0301":
        tokens_per_message = 4  # every message follows <|start|>{role/name}\n{content}<|end|>\n
        tokens_per_name = -1  # if there's a name, the role is omitted
    elif model == "gpt-4-0314":
        tokens_per_message = 3
        tokens_per_name = 1
    else:
        raise NotImplementedError(f"""num_tokens_from_messages() is not implemented for model {model}. See https://github.com/openai/openai-python/blob/main/chatml.md for information on how messages are converted to tokens.""")
    num_tokens = 0
    for message in messages:
        num_tokens += tokens_per_message
        for key, value in message.items():
            num_tokens += len(encoding.encode(value))
            if key == "name":
                num_tokens += tokens_per_name
    num_tokens += 3  # every reply is primed with <|start|>assistant<|message|>
    return num_tokens

def get_model_name() -> str:
  model_name = "GPT-3.5" if config.gpt_model == "gpt-3.5-turbo" else "GPT-4"
  return model_name

def get_max_tokens() -> int:
  max_tokens = 4097 if config.gpt_model == "gpt-3.5-turbo" else 8192
  return max_tokens

# Hey, just FYI if you're going through this source code, this is such a lazy and unorganized solution
# You could probably find a better solution (by literally just copy pasting what the API outputs to you as an error message)
# That's one possible option. But just saying, don't do what I did below. I'll come back and change this in the future, this was a temporary 5 minute fix.
def calculate_completion_tokens(convo_history):
    # this function calculates the remaining tokens needed to reach the maximum number of tokens
    # and returns either the number of tokens to complete, or the maximum number of tokens along
    # with a boolean indicating whether the input parsed_string is already at max tokens.
    
    max_tokens = get_max_tokens() # get the maximum tokens allowed 
    num_tokens = num_tokens_from_messages(convo_history, config.gpt_model)
    
    if num_tokens < (max_tokens - 6): # if the number of tokens is below the maximum number allowed
      remaining_tokens = max_tokens - num_tokens # calculate the number of remaining tokens needed to reach the maximum number
      new_completion_tokens = max(remaining_tokens - 6, 6) # set the new completion tokens
      return new_completion_tokens, False # return the new completion tokens with False indicating that the parsed_string is not yet at max tokens
    else:
      return num_tokens, True # return the number of tokens, maximum number of tokens, and True indicating that the parsed_string is already at max tokens

async def send_message(client, message):
  
  model_name = get_model_name()
  sent_message = await message.channel.send(model_name + " is thinking...")
  
  try:
    author = message.author
    new_completion_tokens = calculate_completion_tokens(config.conversation_history)

    if new_completion_tokens[1]:  # Check for error using the third value in the tuple
      await message.channel.send(
      f"> **Error:** Your input is {new_completion_tokens[0]} tokens, but {model_name} accepts a maximum of {get_max_tokens()} tokens. Please reduce the length of your messages. NOTE: Your input's tokens may be slightly smaller than the max token count. This is normal.")
      config.logger.error(f"@{client.user.name}'s input was {new_completion_tokens[0]} tokens, but {model_name} only accepts a maximum of {get_max_tokens()} tokens")
    
    # Use OpenAI to generate a response to the user's message
    response = openai.ChatCompletion.create(model=config.gpt_model,
                                            messages=config.conversation_history,
                                            temperature=0.55,
                                            max_tokens=new_completion_tokens[0],
                                            n=1,
                                            stop=None)
    
    # Parse the response from OpenAI and format it for the user
    parsed_response = response["choices"][0]["message"]["content"]

    # Please never do this if-else statement in real life, this is absolutely the most troll method I could think of.
    # Wiping the array instead of coming up with an actual solution is honestly stupid
    # "Don't fix what's not broke" - Sun Tzu, probably
    if config.toggle_conversation_history == True:
      config.conversation_history.append({"role": "assistant", "content": parsed_response})
    else:
      config.conversation_history = []

    # Add metadata indicating who the response is for
    response = '> ** Response to ' + str(
      author) + '**:' + '>\n\n' + parsed_response
    
    # Determine whether the response should be sent privately or in the channel
    if config.is_private:
      author = await message.author.create_dm()
    else:
      author = message.channel
    
    # If the response is too long to send in a single message...
    if len(parsed_response) > 1900:
        parts = [parsed_response]  # Initialize parts with the original response
        if "```" in parsed_response:  # If it contains code blocks, split into parts
            parts = parsed_response.split("```")

        for i, part in enumerate(parts):
            # If it's a code block, add the backticks back
            part = "```" + part + "```" if i % 2 == 1 else part
            
            # If the part is too long, split it into chunks and send each chunk as a separate message
            for i in range(0, len(part), 1900):
                chunk = part[i:i + 1900]
                await author.send(chunk) 

    else:
      await author.send(parsed_response)

    # Log to console that response has been sent
    config.logger.info("Response sent to user.\nRESPONSE:  " + parsed_response + "\n")

  # Send an error message to the user if an exception is raised, then log the error
  except Exception as e:
    if "That model is currently overloaded" in str(e):
        await message.channel.send(
      "> **Error:** " + model_name + " is currently overloaded. Please try again later!")
    else:
      await message.channel.send(
      "> **Error:** Something went wrong. Please try again later!")
      
    config.logger.error(e)

  # Delete "GPT-4 is thinking..." message
  await sent_message.delete()
  config.logger.debug("Deleted '" + model_name + " is thinking' message.")